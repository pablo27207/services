#!/usr/bin/env python3
"""
Script para generar presentación del Sistema OOGSJ
para las XII Jornadas Nacionales de Ciencias del Mar
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def crear_presentacion():
    # Crear presentación en formato 4:3
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Colores institucionales (azul marino para oceanografía)
    COLOR_TITULO = RGBColor(0, 51, 102)  # Azul marino
    COLOR_SUBTITULO = RGBColor(0, 102, 153)  # Azul claro
    COLOR_TEXTO = RGBColor(40, 40, 40)  # Gris oscuro
    COLOR_ACENTO = RGBColor(0, 153, 204)  # Celeste

    # ==================== SLIDE 1: PORTADA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Fondo azul suave
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)

    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = "Plataforma Modular para Ingesta, Procesamiento y Visualización de Datos Oceanográficos y Meteorológicos en el Golfo San Jorge"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_TITULO
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Autores
    autores_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
    autores_frame = autores_box.text_frame
    autores_frame.word_wrap = True
    autores_frame.text = "Rosales Pablo¹, De Marziani Carlos¹'², Micheletto Matías², García Franco¹"
    autores_frame.paragraphs[0].font.size = Pt(16)
    autores_frame.paragraphs[0].font.color.rgb = COLOR_TEXTO
    autores_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Instituciones
    inst_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(1.5))
    inst_frame = inst_box.text_frame
    inst_frame.word_wrap = True

    p1 = inst_frame.paragraphs[0]
    p1.text = "¹ Facultad de Ingeniería, Universidad Nacional de la Patagonia San Juan Bosco"
    p1.font.size = Pt(12)
    p1.font.color.rgb = COLOR_TEXTO
    p1.alignment = PP_ALIGN.CENTER

    p2 = inst_frame.add_paragraph()
    p2.text = "² Instituto Multidisciplinario para la Investigación y el Desarrollo Productivo y Social de la Cuenca Golfo San Jorge, CONICET"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXTO
    p2.alignment = PP_ALIGN.CENTER

    # Evento
    evento_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
    evento_frame = evento_box.text_frame
    evento_frame.text = "XII Jornadas Nacionales de Ciencias del Mar"
    evento_frame.paragraphs[0].font.size = Pt(16)
    evento_frame.paragraphs[0].font.italic = True
    evento_frame.paragraphs[0].font.color.rgb = COLOR_SUBTITULO
    evento_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Año
    year_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.5))
    year_frame = year_box.text_frame
    year_frame.text = "2024"
    year_frame.paragraphs[0].font.size = Pt(14)
    year_frame.paragraphs[0].font.color.rgb = COLOR_TEXTO
    year_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ==================== SLIDE 2: ROMA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Red de Observación Marina Argentina (ROMA)"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    p = content.add_paragraph()
    p.text = "Contexto Institucional:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.space_after = Pt(10)

    puntos = [
        "Iniciativa interinstitucional para monitoreo marino de largo plazo",
        "Fortalecimiento de capacidades tecnológicas a nivel nacional",
        "Consolidación de una red de datos oceanográficos interoperables",
        "Articulación entre instituciones académicas y de investigación",
        "Este trabajo contribuye a los objetivos de ROMA en la región del Golfo San Jorge"
    ]

    for punto in puntos:
        p = content.add_paragraph()
        p.text = punto
        p.level = 1
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(8)

    # ==================== SLIDE 3: CONTEXTO ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Contexto: Golfo San Jorge y Economía Azul"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    puntos = [
        "Región de alta productividad biológica en el Atlántico Sur",
        "Zona estratégica para la economía azul: pesca, turismo, hidrocarburos",
        "Ecosistema sensible a cambios ambientales y climáticos",
        "Necesidad de datos continuos para investigación y gestión sostenible",
        "Escasez histórica de datos oceanográficos in situ de largo plazo"
    ]

    for punto in puntos:
        p = content.add_paragraph()
        p.text = punto
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXTO
        p.space_before = Pt(10)

    # ==================== SLIDE 4: PROBLEMÁTICA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Desafíos en el Monitoreo Marino"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    puntos = [
        "Datos oceanográficos fragmentados y discontinuos",
        "Falta de integración eficiente entre múltiples fuentes",
        "Dificultad para estudios de series temporales largas",
        "Limitaciones en el acceso centralizado a información en tiempo real",
        "Necesidad de trazabilidad y control de calidad estandarizado"
    ]

    for punto in puntos:
        p = content.add_paragraph()
        p.text = punto
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXTO
        p.space_before = Pt(10)

    # ==================== SLIDE 5: OBJETIVOS ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Objetivos"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    # Objetivo general
    p = content.add_paragraph()
    p.text = "Objetivo General:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO

    p = content.add_paragraph()
    p.text = "Desarrollar una plataforma modular para la ingesta, almacenamiento, procesamiento y visualización de datos oceanográficos y meteorológicos en el Golfo San Jorge"
    p.font.size = Pt(17)
    p.font.color.rgb = COLOR_TEXTO
    p.space_after = Pt(15)

    # Objetivos específicos
    p = content.add_paragraph()
    p.text = "Objetivos Específicos:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO

    objetivos_esp = [
        "Integrar datos de múltiples plataformas distribuidas",
        "Asegurar trazabilidad y control de calidad según estándares COI",
        "Proporcionar acceso abierto a datos en tiempo real",
        "Respaldar políticas de desarrollo sostenible en zonas costeras"
    ]

    for obj in objetivos_esp:
        p = content.add_paragraph()
        p.text = obj
        p.level = 1
        p.font.size = Pt(17)
        p.font.color.rgb = COLOR_TEXTO

    # ==================== SLIDE 6: PLATAFORMAS ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Plataformas de Monitoreo Integradas"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    plataformas = [
        ("Estaciones Meteorológicas Davis",
         "• Puerto CR y Caleta Córdova (APPCR)\n• Variables: T, P, viento, precipitación, UV\n• Frecuencia: 10 minutos"),

        ("Mareógrafo - Puerto Comodoro Rivadavia",
         "• Sensor: Valeport TideMaster\n• Variable: Nivel del mar\n• Frecuencia: 10 minutos"),

        ("Boya Oceanográfica CIDMAR-2",
         "• Ubicación: 45.877°S, 67.442°W\n• Variables: Olas, corrientes, radiación PAR\n• Frecuencia: Horaria"),

        ("Modelo de Predicción de Mareas",
         "• Servicio Hidrográfico Naval\n• Actualización: Cada 6 horas")
    ]

    for nombre, detalles in plataformas:
        p = content.add_paragraph()
        p.text = nombre
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACENTO

        p = content.add_paragraph()
        p.text = detalles
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(6)

    # Nota sobre fuente
    p = content.add_paragraph()
    p.text = "\n* Datos adquiridos mediante APIs provistas por IADO"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = COLOR_SUBTITULO

    # ==================== SLIDE 7: VARIABLES OCEANOGRÁFICAS ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Variables Oceanográficas Medidas"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    variables_ocean = [
        ("Altura de Olas", "m", "Boya CIDMAR-2"),
        ("Periodo de Olas", "s", "Boya CIDMAR-2"),
        ("Dirección de Olas", "°", "Boya CIDMAR-2"),
        ("Velocidad de Corriente", "m/s", "Boya CIDMAR-2"),
        ("Dirección de Corriente", "°", "Boya CIDMAR-2"),
        ("Radiación PAR", "µmol/m²/s", "Boya CIDMAR-2"),
        ("Nivel del Mar", "m", "Mareógrafo")
    ]

    for var, unidad, fuente in variables_ocean:
        p = content.add_paragraph()
        p.text = f"{var} ({unidad})"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACENTO

        p = content.add_paragraph()
        p.text = f"  Fuente: {fuente}"
        p.level = 1
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(4)

    # ==================== SLIDE 8: VARIABLES METEOROLÓGICAS ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Variables Meteorológicas Medidas"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    # Dividir en dos columnas
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    # Columna izquierda
    vars_left = [
        "Temperatura del aire",
        "Humedad relativa",
        "Presión barométrica",
        "Velocidad del viento",
        "Dirección del viento",
        "Precipitación acumulada",
        "Tasa de precipitación"
    ]

    for var in vars_left:
        p = left_frame.add_paragraph()
        p.text = f"• {var}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(8)

    # Columna derecha
    vars_right = [
        "Punto de rocío",
        "Índice de calor",
        "Sensación térmica",
        "Radiación solar",
        "Índice UV",
        "Evapotranspiración",
        "Cobertura nubosa"
    ]

    for var in vars_right:
        p = right_frame.add_paragraph()
        p.text = f"• {var}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(8)

    # ==================== SLIDE 9: ARQUITECTURA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Arquitectura Modular del Sistema"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    p = content.add_paragraph()
    p.text = "Componentes Principales:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.space_after = Pt(10)

    componentes = [
        "Contenedores Docker: Arquitectura modular y portable",
        "Tareas asíncronas con Celery: Procesamiento distribuido",
        "PostgreSQL: Base de datos con esquema orientado a metadatos",
        "Control de calidad según normativa COI",
        "Frontend web: Dashboard de control y visualización"
    ]

    for comp in componentes:
        p = content.add_paragraph()
        p.text = comp
        p.level = 1
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(10)

    p = content.add_paragraph()
    p.text = "\n✓ Diseño abierto y replicable\n✓ Sistema operativo 24/7\n✓ Actualización automática programada"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_ACENTO

    # ==================== SLIDE 10: FLUJO DE DATOS ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Flujo de Ingesta y Procesamiento"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    flujo = [
        ("1. Adquisición automática", "APIs externas (IADO, WeatherLink, SHN)"),
        ("2. Ingesta con Celery", "Tareas programadas y distribuidas"),
        ("3. Control de calidad", "Validación según estándares COI"),
        ("4. Almacenamiento", "PostgreSQL con metadatos y trazabilidad"),
        ("5. Procesamiento", "Niveles: crudo → calibrado → derivado"),
        ("6. Visualización", "Dashboard web en tiempo real")
    ]

    for paso, desc in flujo:
        p = content.add_paragraph()
        p.text = paso
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACENTO

        p = content.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(8)

    # ==================== SLIDE 11: CONTROL DE CALIDAD ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Control de Calidad y Trazabilidad"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    p = content.add_paragraph()
    p.text = "Quality Flags (según COI/UNESCO):"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.space_after = Pt(10)

    flags = [
        ("Flag 1", "Dato verificado como bueno"),
        ("Flag 2", "Dato probablemente bueno"),
        ("Flag 3", "Dato sospechoso (requiere revisión)"),
        ("Flag 4", "Dato erróneo (descartado)")
    ]

    for flag, desc in flags:
        p = content.add_paragraph()
        p.text = f"{flag}: {desc}"
        p.level = 1
        p.font.size = Pt(17)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(6)

    p = content.add_paragraph()
    p.text = "\nNiveles de Procesamiento:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.space_after = Pt(10)

    niveles = ["Raw (crudo)", "QC-1 (control básico)", "QC-2 (control avanzado)",
               "Interpolated", "Derived (calculado)"]

    for nivel in niveles:
        p = content.add_paragraph()
        p.text = nivel
        p.level = 1
        p.font.size = Pt(17)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(4)

    # ==================== SLIDE 12: VISUALIZACIÓN ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Dashboard de Visualización Web"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    p = content.add_paragraph()
    p.text = "Características del Dashboard:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.space_after = Pt(10)

    caracteristicas = [
        "Gráficos interactivos de series temporales (D3.js)",
        "Mapas dinámicos con ubicación de plataformas (Leaflet + Anime.js)",
        "Visualización de datos en tiempo real",
        "Panel de supervisión del estado del sistema e ingesta",
        "Exploración de atributos técnicos y metadatos",
        "Exportación de datos (CSV, PDF)",
        "Acceso desde cualquier dispositivo (responsive design)"
    ]

    for caract in caracteristicas:
        p = content.add_paragraph()
        p.text = caract
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(6)

    # ==================== SLIDE 13: APLICACIONES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Aplicaciones Científicas y Prácticas"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    aplicaciones = [
        ("Oceanografía Operativa",
         "Monitoreo continuo, predicción de eventos, circulación costera"),

        ("Biología Marina",
         "Correlación con distribución de especies, productividad primaria"),

        ("Economía Azul",
         "Información para pesca, turismo, navegación, gestión sostenible"),

        ("Cambio Climático",
         "Series temporales largas para estudios de tendencias"),

        ("Seguridad Marítima",
         "Prevención de accidentes, planificación de operaciones"),

        ("Educación y Divulgación",
         "Recurso educativo para instituciones académicas")
    ]

    for titulo, desc in aplicaciones:
        p = content.add_paragraph()
        p.text = titulo
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACENTO

        p = content.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(5)

    # ==================== SLIDE 14: ESTADO ACTUAL ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Estado Actual del Sistema"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    p = content.add_paragraph()
    p.text = "Sistema Operativo:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.space_after = Pt(10)

    estado = [
        "✓ Sistema en operación continua 24/7",
        "✓ 5 plataformas de monitoreo integradas",
        "✓ 54 variables diferentes registradas",
        "✓ Base de datos con series temporales desde 2023",
        "✓ Dashboard web accesible públicamente",
        "✓ Actualizaciones automáticas programadas",
        "✓ Datos abiertos para la comunidad científica"
    ]

    for item in estado:
        p = content.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(8)

    # ==================== SLIDE 15: CONCLUSIONES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Conclusiones"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO

    content = slide.placeholders[1].text_frame
    content.clear()

    conclusiones = [
        "La plataforma modular proporciona una infraestructura robusta y replicable para el monitoreo oceanográfico",

        "La integración eficiente de múltiples fuentes de datos permite una visión holística del ambiente marino",

        "El control de calidad según estándares COI asegura trazabilidad y consistencia de las observaciones",

        "El acceso abierto a datos en tiempo real fortalece la investigación y la toma de decisiones",

        "La arquitectura contribuye a los objetivos de ROMA y respalda políticas de desarrollo sostenible en zonas costeras"
    ]

    for i, concl in enumerate(conclusiones, 1):
        p = content.add_paragraph()
        p.text = f"{i}. {concl}"
        p.font.size = Pt(17)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(10)

    # ==================== SLIDE 16: AGRADECIMIENTOS ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Fondo similar a portada
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Agradecimientos"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_TITULO
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Contenido de agradecimientos
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(3.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    agradecimientos = [
        "• Red de Observación Marina Argentina (ROMA)",
        "• Instituto Argentino de Oceanografía (IADO)",
        "• CIDMAR - Centro de Investigación y Desarrollo en Medio Ambiente y Recursos Marinos",
        "• Servicio de Hidrografía Naval Argentina",
        "• Universidad Nacional de la Patagonia San Juan Bosco",
        "• Administración de Puertos del Puerto de Comodoro Rivadavia (APPCR)",
        "• CONICET",
        "• A todos los investigadores y técnicos que contribuyen al monitoreo"
    ]

    for agr in agradecimientos:
        p = content_frame.add_paragraph()
        p.text = agr
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(8)

    # Contacto
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.8))
    contact_frame = contact_box.text_frame
    contact_frame.text = "¿Preguntas?\nContacto: prosales@unpata.edu.ar"
    contact_frame.paragraphs[0].font.size = Pt(16)
    contact_frame.paragraphs[0].font.bold = True
    contact_frame.paragraphs[0].font.color.rgb = COLOR_ACENTO
    contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Guardar presentación
    prs.save('/home/user/services/Presentacion_OOGSJ_XIIJNCM.pptx')
    print("✅ Presentación creada exitosamente: Presentacion_OOGSJ_XIIJNCM.pptx")
    print(f"📊 Total de diapositivas: {len(prs.slides)}")
    print(f"📐 Formato: 4:3 ({prs.slide_width/914400:.1f}\" x {prs.slide_height/914400:.1f}\")")

if __name__ == "__main__":
    crear_presentacion()
